from __future__ import annotations

import json
import math
import os
import re
import tempfile
from collections import Counter
from io import BytesIO
from typing import Dict, List, Tuple

from django.conf import settings
from django.core.files.storage import default_storage
from django.utils import timezone

from apps.ai_services.services import AIService

from .models import KnowledgeDocument


_BLOOM_KEYWORDS: Dict[int, List[str]] = {
    1: ['define', 'list', 'recall', 'identify', 'label', 'name', 'state', 'match'],
    2: ['describe', 'explain', 'summarize', 'paraphrase', 'classify', 'interpret'],
    3: ['apply', 'use', 'demonstrate', 'solve', 'execute', 'implement'],
    4: ['analyze', 'compare', 'contrast', 'differentiate', 'organize', 'examine'],
    5: ['evaluate', 'justify', 'critique', 'argue', 'defend', 'assess'],
    6: ['create', 'design', 'develop', 'construct', 'produce', 'formulate'],
}

_MODALITY_KEYWORDS = {
    'listening': ['listen', 'audio', 'hearing', 'podcast'],
    'speaking': ['speak', 'pronounce', 'presentation', 'speech'],
    'writing': ['write', 'essay', 'compose', 'draft'],
    'math': ['calculate', 'equation', 'solve', 'math', 'formula'],
    'general_knowledge': ['trivia', 'general knowledge', 'facts'],
}

_VALID_MODALITIES = {'reading', 'writing', 'listening', 'speaking', 'math', 'general_knowledge'}
_TRANSIENT_ERRORS = ('timeout', 'temporar', 'connection', 'reset by peer', 'rate limit', 'service unavailable')


def _error(message: str, code: str, retryable: bool = False, details: dict | None = None) -> dict:
    return {
        'error_code': code,
        'message': message,
        'retryable': retryable,
        'details': details or {},
    }


def _extract_error(exc: Exception, code: str, default_message: str, retryable: bool = False, details: dict | None = None) -> dict:
    message = str(exc) or default_message
    lower = message.lower()
    computed_retryable = retryable or any(marker in lower for marker in _TRANSIENT_ERRORS)
    return _error(message=message, code=code, retryable=computed_retryable, details=details)


def can_extract_document(document: KnowledgeDocument) -> dict:
    if document.content_text:
        return {'ok': True, 'code': 'OK_CONTENT_TEXT'}
    if document.source_type == 'url' and not document.source_url:
        return _error('URL source requires source_url.', 'URL_SOURCE_MISSING', retryable=False)
    if document.source_type == 'url':
        return {'ok': True, 'code': 'OK_URL_SOURCE'}
    if not document.file_path:
        return _error('File source requires file_path or content_text.', 'FILE_SOURCE_MISSING', retryable=False)

    ext = document.file_path.rsplit('.', 1)[-1].lower()
    if ext in {'png', 'jpg', 'jpeg', 'bmp', 'tiff'}:
        try:
            from PIL import Image  # noqa: F401
            import pytesseract  # noqa: F401
        except Exception as exc:
            return _extract_error(
                exc,
                code='OCR_DEPENDENCY_MISSING',
                default_message='OCR dependencies are not installed.',
                retryable=False,
                details={'extension': ext},
            )
        return {'ok': True, 'code': 'OK_OCR_CAPABLE'}

    if ext in {'mp4', 'mov', 'avi', 'mkv', 'webm'}:
        if not settings.OPENAI_API_KEY:
            return _error('OPENAI_API_KEY is required for video transcription.', 'VIDEO_TRANSCRIPTION_UNCONFIGURED', retryable=False)
        try:
            from moviepy.editor import VideoFileClip  # noqa: F401
        except Exception as exc:
            return _extract_error(
                exc,
                code='VIDEO_DEPENDENCY_MISSING',
                default_message='Video processing dependencies are not installed.',
                retryable=False,
                details={'extension': ext},
            )
        return {'ok': True, 'code': 'OK_VIDEO_CAPABLE'}

    if ext in {'mp3', 'wav', 'm4a', 'aac', 'ogg'}:
        if not settings.OPENAI_API_KEY:
            return _error('OPENAI_API_KEY is required for audio transcription.', 'AUDIO_TRANSCRIPTION_UNCONFIGURED', retryable=False)
        return {'ok': True, 'code': 'OK_AUDIO_CAPABLE'}

    return {'ok': True, 'code': 'OK_FILE_SOURCE'}


def estimate_tokens(text: str) -> int:
    return max(1, int(len(text) / 4))


def _detect_modality(text: str) -> str:
    lowered = text.lower()
    for modality, keywords in _MODALITY_KEYWORDS.items():
        if any(word in lowered for word in keywords):
            return modality
    return 'reading'


def _cosine_similarity(vec_a: List[float], vec_b: List[float]) -> float:
    if not vec_a or not vec_b:
        return 0.0
    dot = sum(a * b for a, b in zip(vec_a, vec_b))
    norm_a = math.sqrt(sum(a * a for a in vec_a))
    norm_b = math.sqrt(sum(b * b for b in vec_b))
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)


def _ai_classify_bloom(text: str) -> Tuple[int, int, float, str] | None:
    if not settings.OPENAI_API_KEY:
        return None
    prompt = (
        'Return JSON only with keys primary_bloom(1-6), secondary_bloom(1-6), '
        'confidence(0-1), modality(one of reading,writing,listening,speaking,math,general_knowledge).'
    )
    try:
        raw = AIService().chat_completion(
            [
                {'role': 'system', 'content': 'You classify educational text.'},
                {'role': 'user', 'content': f'{prompt}\n\nText:\n{text[:3000]}'},
            ],
            model='gpt-4o-mini',
        )
        start = raw.find('{')
        end = raw.rfind('}')
        if start == -1 or end == -1:
            return None
        parsed = json.loads(raw[start:end + 1])
        primary = int(parsed.get('primary_bloom', 2))
        secondary = int(parsed.get('secondary_bloom', 3))
        confidence = float(parsed.get('confidence', 0.5))
        modality = str(parsed.get('modality', 'reading'))
        if primary < 1 or primary > 6:
            primary = 2
        if secondary < 1 or secondary > 6:
            secondary = 3
        confidence = max(0.0, min(1.0, confidence))
        if modality not in _VALID_MODALITIES:
            modality = 'reading'
        return primary, secondary, round(confidence, 2), modality
    except Exception:
        return None


def classify_bloom_level(text: str) -> Tuple[int, int, float, str]:
    ai_result = _ai_classify_bloom(text)
    if ai_result:
        return ai_result

    lowered = text.lower()
    scores = {}
    for level, keywords in _BLOOM_KEYWORDS.items():
        scores[level] = sum(lowered.count(word) for word in keywords)
    ranked = sorted(scores.items(), key=lambda item: item[1], reverse=True)
    primary_level, primary_score = ranked[0]
    secondary_level, _secondary_score = ranked[1]
    total = sum(scores.values())
    if total == 0:
        primary_level = 2
        secondary_level = 3
        confidence = 0.4
    else:
        confidence = round(primary_score / total, 2)
    modality = _detect_modality(text)
    return primary_level, secondary_level, confidence, modality


def chunk_text(text: str, max_chars: int = 1200, min_chars: int = 600) -> List[str]:
    sentences = re.split(r'(?<=[.!?])\s+|\n+', text.strip())
    chunks: List[str] = []
    buffer: List[str] = []
    current_len = 0
    for sentence in sentences:
        if not sentence:
            continue
        if current_len + len(sentence) > max_chars and current_len >= min_chars:
            chunks.append(' '.join(buffer).strip())
            buffer = [sentence]
            current_len = len(sentence)
            continue
        buffer.append(sentence)
        current_len += len(sentence) + 1
    if buffer:
        chunks.append(' '.join(buffer).strip())
    return [chunk for chunk in chunks if chunk]


def semantic_chunk_text(text: str, max_chars: int = 1400, min_chars: int = 700) -> List[str]:
    sentences = re.split(r'(?<=[.!?])\s+|\n+', text.strip())
    sentences = [sentence.strip() for sentence in sentences if sentence.strip()]
    if len(sentences) < 8 or not settings.OPENAI_API_KEY:
        return chunk_text(text, max_chars=max_chars, min_chars=min_chars)

    try:
        ai = AIService()
        embeddings = [ai.text_embedding(sentence[:800]) for sentence in sentences]
    except Exception:
        return chunk_text(text, max_chars=max_chars, min_chars=min_chars)

    chunks: List[str] = []
    buffer: List[str] = []
    current_len = 0
    previous_embedding = None
    for sentence, embedding in zip(sentences, embeddings):
        similarity = _cosine_similarity(previous_embedding, embedding) if previous_embedding else 1.0
        should_split = current_len >= min_chars and (current_len + len(sentence) > max_chars or similarity < 0.72)
        if should_split:
            chunks.append(' '.join(buffer).strip())
            buffer = [sentence]
            current_len = len(sentence)
        else:
            buffer.append(sentence)
            current_len += len(sentence) + 1
        previous_embedding = embedding

    if buffer:
        chunks.append(' '.join(buffer).strip())
    return [chunk for chunk in chunks if chunk]


def extract_text_for_document(document: KnowledgeDocument) -> str:
    result = extract_text_for_document_with_status(document)
    return result.get('text', '')


def extract_text_for_document_with_status(document: KnowledgeDocument) -> dict:
    capability = can_extract_document(document)
    if not capability.get('ok'):
        return {'text': '', 'error': capability}

    if document.content_text:
        return {'text': document.content_text, 'error': None}
    if document.source_type == 'url' and document.source_url:
        text, error = extract_text_from_url(document.source_url)
        return {'text': text, 'error': error}
    if not document.file_path:
        return {'text': '', 'error': _error('File source requires file_path or content_text.', 'FILE_SOURCE_MISSING', retryable=False)}

    ext = document.file_path.rsplit('.', 1)[-1].lower()
    if ext in {'png', 'jpg', 'jpeg', 'bmp', 'tiff'}:
        text, error = extract_text_from_image(document.file_path)
        return {'text': text, 'error': error}
    if ext in {'pptx', 'ppt'}:
        text, error = extract_text_from_pptx(document.file_path)
        return {'text': text, 'error': error}
    if ext in {'pdf'}:
        text, error = extract_text_from_pdf(document.file_path)
        return {'text': text, 'error': error}
    if ext in {'mp4', 'mov', 'avi', 'mkv', 'webm'} and settings.OPENAI_API_KEY:
        text, error = extract_text_from_video(document.file_path)
        return {'text': text, 'error': error}
    if ext in {'txt', 'md', 'csv', 'json', 'yaml', 'yml', 'html', 'htm'}:
        try:
            with default_storage.open(document.file_path, 'r') as handle:
                return {'text': handle.read(), 'error': None}
        except Exception as exc:
            return {
                'text': '',
                'error': _extract_error(
                    exc,
                    code='TEXT_FILE_READ_FAILED',
                    default_message='Unable to read text document.',
                    retryable=True,
                ),
            }

    if ext in {'mp3', 'wav', 'm4a', 'aac', 'ogg'} and settings.OPENAI_API_KEY:
        try:
            with default_storage.open(document.file_path, 'rb') as handle:
                return {'text': AIService().transcribe_audio(handle), 'error': None}
        except Exception as exc:
            return {
                'text': '',
                'error': _extract_error(
                    exc,
                    code='AUDIO_TRANSCRIPTION_FAILED',
                    default_message='Audio transcription failed.',
                    retryable=True,
                    details={'extension': ext},
                ),
            }

    try:
        with default_storage.open(document.file_path, 'rb') as handle:
            data = handle.read()
    except Exception as exc:
        return {
            'text': '',
            'error': _extract_error(
                exc,
                code='FILE_READ_FAILED',
                default_message='Failed to read source file.',
                retryable=True,
                details={'extension': ext},
            ),
        }
    text = data.decode('utf-8', errors='ignore')
    return {'text': text.strip(), 'error': None}


def extract_text_from_pdf(file_path: str) -> tuple[str, dict | None]:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='PDF_DEPENDENCY_MISSING',
            default_message='PDF extraction dependency missing.',
            retryable=False,
        )
    try:
        with default_storage.open(file_path, 'rb') as handle:
            reader = PdfReader(handle)
            pages = [page.extract_text() or '' for page in reader.pages]
            return '\n'.join(pages).strip(), None
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='PDF_EXTRACTION_FAILED',
            default_message='PDF extraction failed.',
            retryable=True,
        )


def extract_text_from_pptx(file_path: str) -> tuple[str, dict | None]:
    try:
        from pptx import Presentation
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='PPTX_DEPENDENCY_MISSING',
            default_message='PPTX extraction dependency missing.',
            retryable=False,
        )
    try:
        with default_storage.open(file_path, 'rb') as handle:
            data = handle.read()
        presentation = Presentation(BytesIO(data))
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    text_runs.append(shape.text)
        return '\n'.join(text_runs).strip(), None
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='PPTX_EXTRACTION_FAILED',
            default_message='PPTX extraction failed.',
            retryable=True,
        )


def extract_text_from_url(url: str) -> tuple[str, dict | None]:
    try:
        import requests
        from bs4 import BeautifulSoup
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='URL_PARSER_DEPENDENCY_MISSING',
            default_message='URL parsing dependencies are missing.',
            retryable=False,
        )
    try:
        response = requests.get(url, timeout=12)
        response.raise_for_status()
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='URL_FETCH_FAILED',
            default_message='Failed to fetch URL content.',
            retryable=True,
            details={'url': url},
        )
    content_type = response.headers.get('content-type', '').lower()
    body = response.text
    if 'html' in content_type:
        try:
            soup = BeautifulSoup(body, 'html.parser')
            for tag in soup(['script', 'style', 'noscript']):
                tag.decompose()
            text = soup.get_text(separator=' ')
            return re.sub(r'\s+', ' ', text).strip(), None
        except Exception:
            pass
    return re.sub(r'\s+', ' ', body).strip(), None


def extract_text_from_image(file_path: str) -> tuple[str, dict | None]:
    try:
        from PIL import Image
        import pytesseract
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='OCR_DEPENDENCY_MISSING',
            default_message='OCR dependencies are missing.',
            retryable=False,
        )
    try:
        with default_storage.open(file_path, 'rb') as handle:
            image = Image.open(handle)
            return pytesseract.image_to_string(image).strip(), None
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='OCR_EXTRACTION_FAILED',
            default_message='OCR extraction failed.',
            retryable=True,
        )


def extract_text_from_video(file_path: str) -> tuple[str, dict | None]:
    try:
        from moviepy.editor import VideoFileClip
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='VIDEO_DEPENDENCY_MISSING',
            default_message='Video dependencies are missing.',
            retryable=False,
        )
    video_temp = None
    audio_temp = None
    try:
        with default_storage.open(file_path, 'rb') as handle:
            video_data = handle.read()
        video_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.mp4')
        video_temp.write(video_data)
        video_temp.flush()
        video_temp.close()
        audio_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
        audio_temp.close()
        clip = VideoFileClip(video_temp.name)
        if clip.audio is None:
            return '', _error('Video file has no audio track.', 'VIDEO_AUDIO_TRACK_MISSING', retryable=False)
        clip.audio.write_audiofile(audio_temp.name, logger=None)
        clip.close()
        with open(audio_temp.name, 'rb') as audio_file:
            return AIService().transcribe_audio(audio_file).strip(), None
    except Exception as exc:
        return '', _extract_error(
            exc,
            code='VIDEO_TRANSCRIPTION_FAILED',
            default_message='Video transcription failed.',
            retryable=True,
        )
    finally:
        for temp_file in (video_temp, audio_temp):
            if not temp_file:
                continue
            try:
                os.unlink(temp_file.name)
            except Exception:
                pass


def summarize_bloom_distribution(levels: List[int]) -> Dict[str, int]:
    counts = Counter(levels)
    return {str(level): counts.get(level, 0) for level in range(1, 7)}


def infer_audience_profile(document: KnowledgeDocument) -> Dict[str, list]:
    roles = list(
        document.organization.role_competency_mappings.values_list('role_name', flat=True).distinct()
    )
    competencies = list(
        document.organization.competencies.values_list('name', flat=True)[:10]
    )
    return {
        'roles': roles[:10],
        'competencies': competencies,
        'generated_at': timezone.now().isoformat(),
    }


def generate_embedding(text: str) -> List[float] | None:
    if not settings.OPENAI_API_KEY:
        return None
    try:
        return AIService().text_embedding(text[:8000])
    except Exception:
        return None
